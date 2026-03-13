"""Universal file processor — recognizes and extracts content from any file type.

Supported:
  Text:    .txt .md .py .js .ts .html .css .json .yaml .yml .xml .csv .log
  Office:  .pdf .docx .xlsx .xls .pptx .ppt .doc
  Image:   .jpg .jpeg .png .gif .webp .bmp .tiff .svg
  Audio:   .mp3 .wav .ogg .flac .m4a .webm .opus
  Archive: .zip (list contents)
  Code:    any text-based source file
"""

import io
import mimetypes
import os
import zipfile
from pathlib import Path
from typing import Optional

import chardet
import structlog

logger = structlog.get_logger()


class FileInfo:
    def __init__(self):
        self.name: str = ""
        self.size: int = 0
        self.mime_type: str = ""
        self.category: str = "unknown"
        self.content: str = ""
        self.summary: str = ""
        self.error: Optional[str] = None
        self.is_image: bool = False
        self.is_audio: bool = False
        self.pages: int = 0
        self.rows: int = 0
        self.sheets: int = 0
        self.slides: int = 0


def detect_mime(path: Path, data: bytes) -> str:
    try:
        import magic
        return magic.from_buffer(data[:2048], mime=True)
    except Exception:
        mime, _ = mimetypes.guess_type(str(path))
        return mime or "application/octet-stream"


def extract_text_file(data: bytes) -> str:
    detected = chardet.detect(data)
    enc = detected.get("encoding") or "utf-8"
    try:
        return data.decode(enc, errors="replace")
    except Exception:
        return data.decode("utf-8", errors="replace")


def extract_pdf(data: bytes) -> tuple[str, int]:
    try:
        from pypdf import PdfReader
        reader = PdfReader(io.BytesIO(data))
        pages = len(reader.pages)
        text = "\n\n".join(
            p.extract_text() or "" for p in reader.pages[:50]
        )
        return text.strip(), pages
    except Exception as e:
        return f"[PDF extraction error: {e}]", 0


def extract_docx(data: bytes) -> str:
    try:
        from docx import Document
        doc = Document(io.BytesIO(data))
        return "\n".join(p.text for p in doc.paragraphs if p.text.strip())
    except Exception as e:
        return f"[DOCX extraction error: {e}]"


def extract_xlsx(data: bytes) -> tuple[str, int, int]:
    try:
        import openpyxl
        wb = openpyxl.load_workbook(io.BytesIO(data), read_only=True, data_only=True)
        sheets = len(wb.sheetnames)
        rows_total = 0
        parts = []
        for name in wb.sheetnames[:5]:
            ws = wb[name]
            rows = []
            for i, row in enumerate(ws.iter_rows(values_only=True)):
                if i > 200:
                    rows.append(f"... (truncated, {ws.max_row} rows total)")
                    break
                cells = [str(c) if c is not None else "" for c in row]
                rows.append("\t".join(cells))
            rows_total += ws.max_row or 0
            parts.append(f"=== Sheet: {name} ===\n" + "\n".join(rows))
        return "\n\n".join(parts), sheets, rows_total
    except Exception as e:
        return f"[XLSX extraction error: {e}]", 0, 0


def extract_pptx(data: bytes) -> tuple[str, int]:
    try:
        from pptx import Presentation
        prs = Presentation(io.BytesIO(data))
        slides = len(prs.slides)
        parts = []
        for i, slide in enumerate(prs.slides):
            texts = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    texts.append(shape.text.strip())
            if texts:
                parts.append(f"[Slide {i+1}]\n" + "\n".join(texts))
        return "\n\n".join(parts), slides
    except Exception as e:
        return f"[PPTX extraction error: {e}]", 0


def extract_csv(data: bytes) -> tuple[str, int]:
    try:
        import pandas as pd
        enc = chardet.detect(data).get("encoding") or "utf-8"
        df = pd.read_csv(io.BytesIO(data), encoding=enc, nrows=500,
                         on_bad_lines="skip")
        rows = len(df)
        preview = df.head(20).to_string(index=False)
        summary = f"Shape: {rows} rows × {len(df.columns)} columns\nColumns: {', '.join(df.columns.tolist())}\n\n{preview}"
        return summary, rows
    except Exception as e:
        return f"[CSV extraction error: {e}]", 0


def extract_zip_listing(data: bytes) -> str:
    try:
        with zipfile.ZipFile(io.BytesIO(data)) as zf:
            names = zf.namelist()
            listing = "\n".join(names[:100])
            if len(names) > 100:
                listing += f"\n... and {len(names)-100} more files"
            return f"ZIP Archive — {len(names)} files:\n{listing}"
    except Exception as e:
        return f"[ZIP listing error: {e}]"


TEXT_EXTENSIONS = {
    ".txt", ".md", ".py", ".js", ".ts", ".jsx", ".tsx", ".html", ".htm",
    ".css", ".scss", ".json", ".yaml", ".yml", ".xml", ".toml", ".ini",
    ".cfg", ".conf", ".sh", ".bash", ".zsh", ".env", ".gitignore",
    ".sql", ".graphql", ".rs", ".go", ".java", ".cpp", ".c", ".h",
    ".php", ".rb", ".swift", ".kt", ".log", ".rst",
}


async def process_file(filename: str, data: bytes) -> FileInfo:
    """Main entry point — detect type and extract content."""
    info = FileInfo()
    info.name = filename
    info.size = len(data)

    path = Path(filename)
    ext = path.suffix.lower()
    info.mime_type = detect_mime(path, data)
    mime = info.mime_type

    logger.info("file_processing", name=filename, size=info.size, mime=mime)

    try:
        # ── Images ──────────────────────────────────────────────────────────
        if mime.startswith("image/") or ext in {".jpg",".jpeg",".png",".gif",".webp",".bmp",".tiff",".svg"}:
            info.category = "image"
            info.is_image = True
            if ext == ".svg":
                info.content = extract_text_file(data)[:2000]
                info.summary = f"SVG vector image ({info.size} bytes)"
            else:
                try:
                    from PIL import Image
                    img = Image.open(io.BytesIO(data))
                    w, h = img.size
                    mode = img.mode
                    info.summary = f"Image: {w}×{h}px, mode={mode}, format={img.format}"
                    info.content = info.summary
                except Exception as e:
                    info.summary = f"Image file ({mime})"
                    info.content = info.summary

        # ── Audio ────────────────────────────────────────────────────────────
        elif mime.startswith("audio/") or ext in {".mp3",".wav",".ogg",".flac",".m4a",".webm",".opus",".aac"}:
            info.category = "audio"
            info.is_audio = True
            info.summary = f"Audio file: {filename} ({info.size/1024:.1f} KB)"
            info.content = info.summary

        # ── PDF ──────────────────────────────────────────────────────────────
        elif mime == "application/pdf" or ext == ".pdf":
            info.category = "document"
            text, pages = extract_pdf(data)
            info.content = text[:8000]
            info.pages = pages
            info.summary = f"PDF: {pages} pages, {len(text)} chars extracted"

        # ── DOCX ─────────────────────────────────────────────────────────────
        elif ext in {".docx", ".doc"} or "wordprocessingml" in mime:
            info.category = "document"
            text = extract_docx(data)
            info.content = text[:8000]
            info.summary = f"Word document: {len(text)} chars"

        # ── XLSX ─────────────────────────────────────────────────────────────
        elif ext in {".xlsx", ".xls"} or "spreadsheet" in mime or "excel" in mime:
            info.category = "spreadsheet"
            text, sheets, rows = extract_xlsx(data)
            info.content = text[:8000]
            info.sheets = sheets
            info.rows = rows
            info.summary = f"Excel: {sheets} sheets, ~{rows} rows"

        # ── PPTX ─────────────────────────────────────────────────────────────
        elif ext in {".pptx", ".ppt"} or "presentationml" in mime:
            info.category = "presentation"
            text, slides = extract_pptx(data)
            info.content = text[:8000]
            info.slides = slides
            info.summary = f"PowerPoint: {slides} slides"

        # ── CSV ──────────────────────────────────────────────────────────────
        elif ext == ".csv" or mime == "text/csv":
            info.category = "data"
            text, rows = extract_csv(data)
            info.content = text[:8000]
            info.rows = rows
            info.summary = f"CSV: ~{rows} rows"

        # ── ZIP ──────────────────────────────────────────────────────────────
        elif ext == ".zip" or mime == "application/zip":
            info.category = "archive"
            info.content = extract_zip_listing(data)
            info.summary = f"ZIP archive"

        # ── Text / Code ──────────────────────────────────────────────────────
        elif ext in TEXT_EXTENSIONS or mime.startswith("text/"):
            info.category = "text"
            text = extract_text_file(data)
            info.content = text[:8000]
            info.summary = f"Text file: {len(text)} chars, {text.count(chr(10))} lines"

        else:
            info.category = "binary"
            info.content = f"Binary file ({mime}), {info.size} bytes — cannot extract text"
            info.summary = info.content

    except Exception as e:
        logger.error("file_processing_error", name=filename, error=str(e))
        info.error = str(e)
        info.content = f"Error processing file: {e}"
        info.summary = info.content

    return info
